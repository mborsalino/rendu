# Copyright (c) 2025 Mattia Borsalino

import pptx
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches
from PIL import Image

class PptSlideDeck(object):
    """
    Wrapper class for a presentation of specific size and aesthetics.

    This is a pilot project, never completed, that led to the conceival of
    Rendu's HtmlSlidDeck. It is kept for refeence and should therefore not
    be used.
    """

    def __init__(self):

        # attributes
        self.width = 16
        self.height = 9
        self.pad = 0.5

        self.tslide_tbar_height = 0.5
        self.tslide_tbar_left = 0
        self.tslide_tbar_top = self.height/1.5
        self.fslide_tbar_height = 0.3
        self.fslide_tbar_left = 0
        self.fslide_tbar_top = 0.5

        self.tslide = None     # title slide
        self.cslides = []      # content slides
        self.theme_color = RGBColor(255, 0, 0) # red, subclass for other colors
        self.logo = None

        # create ppt
        self.p = pptx.Presentation()
        self.p.slide_width = Inches(self.width)
        self.p.slide_height = Inches(self.height)


    def add_title_slide(self, title):
        """
        Add title slide with the given title.
        """
        if self.tslide is not None:
            raise Exception('Only one title slide is allowed')

        self.tslide = self.p.slides.add_slide(self.p.slide_layouts[6])
        title_bar = self.tslide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(self.tslide_tbar_left),   # x of top left corner ('left')
                                         Inches(self.tslide_tbar_top),    # y of top left corner ('top')
                                         Inches(self.width),              # width
                                         Inches(self.tslide_tbar_height)) # height
        title_bar.shadow.inherit = False
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.theme_color
        title_bar.text = title
        border = title_bar.line
        border.color.rgb = self.theme_color


    def add_figure_slide(self, title, fig_path, fit_to_width=False):
        """
        Add a slide containing the given figure.

        Inputs:
        -------
        title: str
            The title of the slide
        figure: str
            The path of a png or jpeg figure
        """
        s = self.p.slides.add_slide(self.p.slide_layouts[6])
        self.cslides.append(s)

        title_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(self.fslide_tbar_left),
                                       Inches(self.fslide_tbar_top),
                                       Inches(self.width),
                                       Inches(self.fslide_tbar_height))
        title_bar.shadow.inherit = False
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.theme_color
        title_bar.text = title
        border = title_bar.line
        border.color.rgb = self.theme_color

        l, t, w, h = self._compute_centered_pic_size_and_pos(fig_path)
        pic = s.shapes.add_picture(fig_path,
                                   Inches(l),
                                   Inches(t),
                                   Inches(w),
                                   Inches(h))


    def _compute_centered_pic_size_and_pos(self, pic_path, pad=None):
        """
        Get picsize so that it's centered on the available area of the slide

        Return:
        (left, top, width, height)
        """
        if pad is None:
            pad = self.pad

        pic = Image.open(pic_path)
        pic_w, pic_h = pic.size
        pic_ar = pic_w / pic_h
        slide_ar = self.width / self.height

        avail_height = self.height - self.fslide_tbar_top - self.fslide_tbar_height
        if slide_ar > pic_ar: # we need to fit height
            fit_height = avail_height - 2 * pad
            fit_width = fit_height * pic_ar

        else: # we need to fit width
            fit_width = self.width - 2 * pad
            fit_height = fit_width / pic_ar

        left = pad
        top = self.fslide_tbar_top + self.fslide_tbar_height + (avail_height - fit_height) / 2
        return (left, top, fit_width, fit_height)


    def save(self, fname):
        """
        Save pptx to given file name. File name must include '.pptx'
        """
        if not fname.endswith('.pptx'):
            raise Exception('File name must have suffix ".pptx"')

        if len(self.p.slides) == 0 and self.tslide is None:
            raise Exception('Presentation is empty')

        self.p.save(fname)
